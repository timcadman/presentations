import os
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.enum.shapes import PP_PLACEHOLDER
from PIL import Image
from lxml import etree
HERE=os.path.dirname(os.path.abspath(__file__))
OUT=os.path.join(HERE,"molgenis-template-STARTER.pptx")
LOGO=os.path.join(HERE,"molgenis-armadillo","public","molgenis-logo.png")
BLUE="4285F4"; BLACK="000000"
TITLEF="Bebas Neue"; SUBF="IBM Plex Mono"; BODYF="Nunito"
BODY_SZ=14   # single body font size used across ALL content layouts (uniform, no per-slide autofit)
def sub(p,t,**a):
    e=etree.SubElement(p,qn(t));  [e.set(k,v) for k,v in a.items()];  return e
def lst(ph):
    txb=ph._element.find(qn('p:txBody')); l=txb.find(qn('a:lstStyle'))
    if l is None: l=etree.Element(qn('a:lstStyle')); txb.find(qn('a:bodyPr')).addnext(l)
    return l
def set_lvl(ph,level,sz,color,font,bullet=None,marL=None,indent=None,sb=None,sa=None,align=None,italic=False):
    ls=lst(ph); tag=qn(f'a:lvl{level}pPr'); lvl=ls.find(tag)
    if lvl is None: lvl=etree.SubElement(ls,tag)
    for c in list(lvl): lvl.remove(c)
    if align: lvl.set('algn',align)
    if marL is not None: lvl.set('marL',str(marL))
    if indent is not None: lvl.set('indent',str(indent))
    if sb is not None: sub(sub(lvl,'a:spcBef'),'a:spcPts',val=str(sb))
    if sa is not None: sub(sub(lvl,'a:spcAft'),'a:spcPts',val=str(sa))
    if bullet=='none': sub(lvl,'a:buNone')
    elif bullet=='dot': sub(lvl,'a:buFont',typeface='Arial'); sub(lvl,'a:buChar',char='•')
    d=sub(lvl,'a:defRPr',sz=str(int(sz*100)))
    if italic: d.set('i','1')
    f=sub(d,'a:solidFill'); sub(f,'a:srgbClr',val=color); sub(d,'a:latin',typeface=font)
def title_block(ph,head_sz):                      # lvl0 heading, lvl1 subheading
    ph._element.find(qn('p:txBody')).find(qn('a:bodyPr')).set('anchor','t')   # top-anchored
    set_lvl(ph,1,head_sz,BLUE,TITLEF,align='l',bullet='none')
    set_lvl(ph,2,18,BLUE,SUBF,align='l',bullet='none',sb=400)
def body_block(ph,size=BODY_SZ):
    ph._element.find(qn('p:txBody')).find(qn('a:bodyPr')).set('anchor','t')  # top-anchored: bullets grow down, never overlap title
    set_lvl(ph,1,size,BLACK,BODYF,bullet='dot',marL=228600,indent=-228600,sa=500)
    set_lvl(ph,2,size-2,BLACK,BODYF,bullet='dot',marL=685800,indent=-228600)
def pos(ph,x,y,w,h): ph.left=Inches(x); ph.top=Inches(y); ph.width=Inches(w); ph.height=Inches(h)
def to_pic(ph): ph._element.find('.//'+qn('p:ph')).set('type','pic')
def nid(s):
    ids=[int(e.get('id')) for e in s.iter(qn('p:cNvPr')) if e.get('id','').isdigit()]
    return (max(ids)+1) if ids else 100
def lname(prs,n):
    for l in prs.slide_layouts:
        if l.name==n: return l
def phd(l): return {p.placeholder_format.idx:p for p in l.placeholders}
def byt(l,t): return [p for p in l.placeholders if p.placeholder_format.type==t]
def add_bar(layout,x,y,w):
    sp=layout.shapes._spTree; s=sp.add_autoshape(nid(sp),"Accent","rect",Inches(x),Inches(y),Inches(w),Inches(0.05))
    st=s.find(qn('p:style'));  s.remove(st) if st is not None else None
    spPr=s.find(qn('p:spPr'))
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill'):
        e=spPr.find(qn(tag));  spPr.remove(e) if e is not None else None
    f=sub(spPr,'a:solidFill'); sub(f,'a:srgbClr',val=BLUE)
    ln=spPr.find(qn('a:ln'));  ln=sub(spPr,'a:ln') if ln is None else ln
    for c in list(ln): ln.remove(c)
    sub(ln,'a:noFill')
def add_logo(container,part,x,y,gw,gh):
    container._spTree.add_pic(nid(container._spTree),"Logo","Logo",part.get_or_add_image_part(LOGO)[1],int(x),int(y),gw,gh)

prs=Presentation(); SW,SH=10.0,5.625
prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)   # default template is 4:3; force 16:9
for s in list(prs.slides._sldIdLst):
    prs.part.drop_rel(s.get(qn('r:id'))); prs.slides._sldIdLst.remove(s)
lw,lh=Image.open(LOGO).size; gw=Inches(1.0); gh=int(gw*lh/lw)
LOGO_BR=(int(Inches(SW)-gw-Inches(0.3)),int(Inches(SH)-gh-Inches(0.2)))
LOGO_BL=(int(Inches(0.4)),int(Inches(SH)-gh-Inches(0.2)))

# Title
L=lname(prs,"Title Slide"); L._element.cSld.set('name',"Title")
ct=byt(L,PP_PLACEHOLDER.CENTER_TITLE)[0]; st=byt(L,PP_PLACEHOLDER.SUBTITLE)[0]
add_bar(L,0.7,0.95,7.0)                                   # bar above the title block
# heading (lvl0) + subheading (lvl1) both flow inside the title placeholder -> no overlap when heading wraps
pos(ct,0.65,1.12,8.7,2.3); ct._element.find(qn('p:txBody')).find(qn('a:bodyPr')).set('anchor','t')
set_lvl(ct,1,55,BLUE,TITLEF,align='l')
set_lvl(ct,2,25,BLUE,SUBF,bullet='none',align='l',sb=300)
st._element.getparent().remove(st._element)              # drop unused separate subtitle placeholder
# name/role are added per-slide by the builder (so they're editable on the slide, not the layout)

# Bullets
L=lname(prs,"Title and Content"); L._element.cSld.set('name',"Bullets")
add_bar(L,0.7,0.6,3.0); t=byt(L,PP_PLACEHOLDER.TITLE)[0]; o=byt(L,PP_PLACEHOLDER.OBJECT)[0]
pos(t,0.7,0.9,8.6,1.4); title_block(t,40); pos(o,0.7,2.35,8.6,2.4); body_block(o)

# Photo Right (full-bleed image right, text left)
L=lname(prs,"Content with Caption"); L._element.cSld.set('name',"Photo Right")
add_bar(L,0.7,0.6,2.8); p=phd(L)
pos(p[0],0.7,0.9,5.0,1.4); title_block(p[0],40)
pos(p[2],0.7,2.35,5.0,2.4); body_block(p[2])
to_pic(p[1]); pos(p[1],6.0,0,4.0,SH)
add_logo(L.shapes,L.part,*LOGO_BL,gw,gh)

# Photo Left (full-bleed image left, text right)
L=lname(prs,"Two Content"); L._element.cSld.set('name',"Photo Left")
p=phd(L); to_pic(p[1]); pos(p[1],0,0,4.0,SH)
add_bar(L,4.6,0.6,2.8)
pos(p[0],4.6,0.9,4.8,1.4); title_block(p[0],40)
pos(p[2],4.6,2.35,4.8,2.4); body_block(p[2])

# Image Left (contained image left, bullets right)
L=lname(prs,"Picture with Caption"); L._element.cSld.set('name',"Image Left")
add_bar(L,0.7,0.6,2.8); p=phd(L)
pos(p[0],0.7,0.9,8.6,1.4); title_block(p[0],40)     # title spans top
pos(p[1],0.7,2.35,4.2,2.5)                           # PICTURE left (contained region)
pos(p[2],5.2,2.35,4.2,2.5); body_block(p[2])        # bullets right

# Image Right (contained image right, bullets left)  -- from Comparison
L=lname(prs,"Comparison"); L._element.cSld.set('name',"Image Right")
add_bar(L,0.7,0.6,2.8); p=phd(L)
# keep title(0), one body for bullets, one obj->pic; remove extras
pos(p[0],0.7,0.9,8.6,1.4); title_block(p[0],40)
pos(p[1],0.7,2.35,4.2,2.5); body_block(p[1])        # bullets left (BODY idx1)
to_pic(p[2]); pos(p[2],5.2,2.35,4.2,2.5)            # picture right (OBJECT idx2)
# remove the spare placeholders idx3, idx4
for idx in (3,4):
    if idx in p:
        p[idx]._element.getparent().remove(p[idx]._element)

# Table
L=lname(prs,"Title Only"); L._element.cSld.set('name',"Table")
add_bar(L,0.7,0.6,3.0); t=byt(L,PP_PLACEHOLDER.TITLE)[0]; pos(t,0.7,0.9,8.6,1.4); title_block(t,40)

for m in prs.slide_masters: add_logo(m.shapes,m.part,*LOGO_BR,gw,gh)
keep={"Title","Bullets","Photo Right","Photo Left","Image Left","Image Right","Table"}
for m in prs.slide_masters:
    l=m.element.get_or_add_sldLayoutIdLst()
    for sid in list(l):
        rid=sid.get(qn('r:id'))
        if m.part.related_part(rid)._element.cSld.get('name') not in keep:
            l.remove(sid); m.part.drop_rel(rid)
prs.save(OUT); print("7-layout template:", [x.name for x in prs.slide_layouts])
