"""
build_deck.py — generate an editable, branded MOLGENIS PowerPoint from slide specs.

All styling (fonts, colours, accent bars, logo, layout geometry) lives in the
template (molgenis-template-STARTER.pptx). This builder only routes content into
the template's layouts and draws the two card-based slides (full-stack, demo).

Usage:
    python3 build_deck.py            # builds /tmp/molgenis-demo.pptx from DEMO_DECK

Each slide is a dict with a "type":
  {"type":"title",   "heading":..., "subheading":..., "name":..., "role":...}
  {"type":"bullets", "heading":..., "subheading":..., "bullets":[...]}
  {"type":"photo",   "side":"right"|"left", "heading":..., "subheading":..., "bullets":[...], "image":path}
  {"type":"image",   "side":"right"|"left", "heading":..., "subheading":..., "bullets":[...], "image":path}
  {"type":"table",   "heading":..., "subheading":..., "rows":[[...],...]}
  {"type":"stack",   "heading":..., "subheading":..., "rows":[(label,desc),...]}
  {"type":"demo",    "step":..., "image":path, "text":...}
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from PIL import Image as PImage
import os

TEMPLATE = os.path.join(os.path.dirname(__file__), "molgenis-template-STARTER.pptx")
ASSET_ROOT = os.path.dirname(__file__)
BLACK = RGBColor(0x22, 0x22, 0x22)
# FAIR gradient (from DemoSlide.vue): label, fill hex, text hex
GRAD = [("Local data", "93c5fd", "ffffff"), ("Catalogue", "60a5fa", "ffffff"),
        ("Request", "3b82f6", "ffffff"), ("Access", "2563eb", "ffffff"),
        ("Analyse", "1e3a5f", "ffffff")]


def _sub(parent, tag, **a):
    e = etree.SubElement(parent, qn(tag))
    for k, v in a.items():
        e.set(k, v)
    return e


class DeckBuilder:
    def __init__(self, template=TEMPLATE):
        self.prs = Presentation(template)

    # ---- helpers ----
    def _layout(self, name):
        for l in self.prs.slide_layouts:
            if l.name == name:
                return l
        raise KeyError(f"layout {name!r} not in template")

    @staticmethod
    def _ph(slide, ptype):
        for p in slide.placeholders:
            if p.placeholder_format.type == ptype:
                return p
        return None

    def _title_sub(self, slide, heading, subheading):
        ph = self._ph(slide, PP_PLACEHOLDER.TITLE) or self._ph(slide, PP_PLACEHOLDER.CENTER_TITLE)
        tf = ph.text_frame
        tf.clear()
        tf.paragraphs[0].add_run().text = heading or ""
        if subheading:
            p = tf.add_paragraph(); p.level = 1; p.add_run().text = subheading

    def _bullets(self, ph, items):
        tf = ph.text_frame; tf.clear()
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = 0; p.add_run().text = it

    def _fit_picture(self, slide, ph, img, border=True, scale=1.0, valign="top"):
        """Contain-fit an image inside the placeholder's rect (no crop).
        scale<1 shrinks within the rect; valign 'top' keeps the fixed top, 'middle' centres."""
        L, T, W, H = ph.left, ph.top, ph.width, ph.height
        ph._element.getparent().remove(ph._element)
        iw, ih = PImage.open(img).size
        ar, AR = iw / ih, W / H
        if ar > AR:
            w = int(W * scale); h = int(W * scale / ar)
        else:
            h = int(H * scale); w = int(H * scale * ar)
        x = int(L + (W - w) / 2)
        y = int(T) if valign == "top" else int(T + (H - h) / 2)
        pic = slide.shapes.add_picture(img, x, y, w, h)
        if border:
            pic.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9); pic.line.width = Pt(1)
        return pic

    # ---- slide types ----
    def title(self, heading, subheading="", name="Tim Cadman", role="Senior Data Scientist", **_):
        s = self.prs.slides.add_slide(self._layout("Title"))
        self._title_sub(s, heading, subheading)
        # editable name/role text box bottom-left
        tb = s.shapes.add_textbox(Inches(0.7), self.prs.slide_height - Inches(1.05), Inches(4), Inches(0.8))
        tf = tb.text_frame
        r = tf.paragraphs[0].add_run(); r.text = name
        r.font.bold = True; r.font.size = Pt(14); r.font.name = "Nunito"; r.font.color.rgb = BLACK
        if role:
            p = tf.add_paragraph(); r = p.add_run(); r.text = role
            r.font.size = Pt(11); r.font.name = "Nunito"; r.font.color.rgb = BLACK
        return s

    def bullets(self, heading, subheading="", bullets=(), **_):
        s = self.prs.slides.add_slide(self._layout("Bullets"))
        self._title_sub(s, heading, subheading)
        self._bullets(self._ph(s, PP_PLACEHOLDER.OBJECT), bullets)
        return s

    def photo(self, side, heading, subheading="", bullets=(), image=None, **_):
        s = self.prs.slides.add_slide(self._layout(f"Photo {side.capitalize()}"))
        self._title_sub(s, heading, subheading)
        self._bullets(self._ph(s, PP_PLACEHOLDER.BODY) or self._ph(s, PP_PLACEHOLDER.OBJECT), bullets)
        if image:
            self._ph(s, PP_PLACEHOLDER.PICTURE).insert_picture(image)  # cover/crop = full-bleed
        return s

    def image(self, side, heading, subheading="", bullets=(), image=None, border=True, **_):
        s = self.prs.slides.add_slide(self._layout(f"Image {side.capitalize()}"))
        self._title_sub(s, heading, subheading)
        # screenshot is top-aligned in its region; bullets sit at the layout's fixed top.
        # both share the same top point and the layout height never changes per slide.
        if image:
            # logos (border=False) read better smaller and vertically centred
            kw = dict(border=False, scale=0.62, valign="middle") if not border else {}
            self._fit_picture(s, self._ph(s, PP_PLACEHOLDER.PICTURE), image, **kw)
        self._bullets(self._ph(s, PP_PLACEHOLDER.BODY), bullets)
        return s

    def two_images(self, heading, subheading="", images=(), **_):
        s = self.prs.slides.add_slide(self._layout("Bullets"))
        self._title_sub(s, heading, subheading)
        ph = self._ph(s, PP_PLACEHOLDER.OBJECT) or self._ph(s, PP_PLACEHOLDER.BODY)
        if ph is not None:
            ph._element.getparent().remove(ph._element)
        top, gap, margin = Inches(2.0), Inches(0.4), Inches(0.7)
        bottom = self.prs.slide_height - Inches(0.6)
        H = bottom - top
        halfW = int((self.prs.slide_width - margin * 2 - gap) / 2)
        imgs = [i for i in images if i][:2]
        for i, img in enumerate(imgs):
            iw, ih = PImage.open(img).size
            ar, AR = iw / ih, halfW / H
            if ar > AR:
                w = halfW; h = int(halfW / ar)
            else:
                h = int(H); w = int(H * ar)
            x = int(margin) + i * (halfW + int(gap)) + (halfW - w) // 2
            pic = s.shapes.add_picture(img, x, int(top), w, h)
            pic.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9); pic.line.width = Pt(1)
        return s

    def slide_image(self, image=None, **_):
        """Whole-slide image fallback for slides that can't be made editable
        (e.g. a bespoke Vue diagram). The image already contains its own title."""
        s = self.prs.slides.add_slide(self._layout("Bullets"))
        for p in list(s.placeholders):
            p._element.getparent().remove(p._element)
        if image:
            s.shapes.add_picture(image, 0, 0, self.prs.slide_width, self.prs.slide_height)
        return s

    def table(self, heading, subheading="", rows=(), **_):
        s = self.prs.slides.add_slide(self._layout("Table"))
        self._title_sub(s, heading, subheading)
        r, c = len(rows), len(rows[0])
        t = s.shapes.add_table(r, c, Inches(0.7), Inches(2.35), Inches(8.6), Inches(2.5)).table
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                t.cell(ri, ci).text = str(val)
        return s

    def stack(self, heading, subheading="", rows=None, **_):
        rows = rows or [(g[0], "") for g in GRAD]
        s = self.prs.slides.add_slide(self._layout("Table"))
        self._title_sub(s, heading, subheading)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.15), Inches(5.2), Inches(2.95))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xEE, 0xF1, 0xF4)
        box.line.fill.background(); box.shadow.inherit = False
        cy = 2.32
        for (label, desc), (_, hexv, _t) in zip(rows, GRAD):
            self._card(s, 0.95, cy, 4.7, 0.48, hexv, label, RGBColor(0xff, 0xff, 0xff), 100000, desc)
            cy += 0.54
        return s

    def demo(self, step, image=None, text="", subheading=None, **_):
        s = self.prs.slides.add_slide(self._layout("Table"))
        self._title_sub(s, "Demo", subheading if subheading is not None else step)
        tb = s.shapes.add_textbox(Inches(0.7), Inches(2.05), Inches(8.5), Inches(0.5))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = text
        r.font.name = "Nunito"; r.font.size = Pt(15); r.font.color.rgb = BLACK
        fx, fy, fw, fh = 0.7, 2.65, 5.6, 2.5
        frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(fx), Inches(fy), Inches(fw), Inches(fh))
        frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF7); frame.shadow.inherit = False
        frame.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); frame.line.width = Pt(1.5)
        _sub(frame._element.spPr.find(qn('a:ln')), 'a:prstDash').set('val', 'dash')
        if image:
            iw, ih = PImage.open(image).size; ar, AR = iw / ih, fw / fh
            if ar > AR: w = fw - 0.2; h = w / ar
            else: h = fh - 0.2; w = h * ar
            s.shapes.add_picture(image, Inches(fx + (fw - w) / 2), Inches(fy + (fh - h) / 2), Inches(w), Inches(h))
        cy = 2.65
        for label, hexv, txc in GRAD:
            active = (label == step)
            self._card(s, 6.5, cy, 2.8, 0.42, hexv, label, RGBColor.from_string(txc),
                       100000 if active else 25000, "", active)
            cy += 0.47
        return s

    def _card(self, s, x, y, w, h, hexv, label, txt_col, alpha, desc="", active=True):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        try: sh.adjustments[0] = 0.2
        except Exception: pass
        sh.shadow.inherit = False
        spPr = sh._element.spPr
        st = sh._element.find(qn('p:style'));  sh._element.remove(st) if st is not None else None
        for tag in ('a:solidFill', 'a:noFill'):
            e = spPr.find(qn(tag));  spPr.remove(e) if e is not None else None
        f = etree.Element(qn('a:solidFill')); c = _sub(f, 'a:srgbClr', val=hexv)
        if alpha < 100000: _sub(c, 'a:alpha', val=str(alpha))
        spPr.find(qn('a:prstGeom')).addnext(f)
        ln = spPr.find(qn('a:ln')) or _sub(spPr, 'a:ln')
        for e in list(ln): ln.remove(e)
        _sub(ln, 'a:noFill')
        lt = s.shapes.add_textbox(Inches(x + 0.18), Inches(y), Inches(w * 0.55), Inches(h))
        lt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = lt.text_frame.paragraphs[0].add_run(); r.text = label
        r.font.bold = True; r.font.size = Pt(12); r.font.name = "Nunito"
        r.font.color.rgb = txt_col                                    # correct XML placement
        if alpha < 100000:
            srgb = r._r.find(qn('a:rPr')).find(qn('a:solidFill')).find(qn('a:srgbClr'))
            _sub(srgb, 'a:alpha', val='35000')
        if desc:
            dt = s.shapes.add_textbox(Inches(x + w * 0.4), Inches(y), Inches(w * 0.57), Inches(h))
            dt.text_frame.word_wrap = True; dt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = dt.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = desc; r.font.italic = True; r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(0xff, 0xff, 0xff); r.font.name = "Nunito"

    def build(self, deck):
        for spec in deck:
            t = spec["type"]
            if t == "title": self.title(**spec)
            elif t == "bullets": self.bullets(**spec)
            elif t == "photo": self.photo(**spec)
            elif t == "image": self.image(**spec)
            elif t == "two_images": self.two_images(**spec)
            elif t == "slide_image": self.slide_image(**spec)
            elif t == "table": self.table(**spec)
            elif t == "stack": self.stack(**spec)
            elif t == "demo": self.demo(**spec)
            else: raise ValueError(f"unknown slide type {t!r}")
        return self.prs


def _p(rel):  # resolve a public asset path
    return os.path.join(ASSET_ROOT, "molgenis-demonstrator", rel)


DEMO_DECK = [
    {"type": "title", "heading": "From Data Discovery to Federated Analysis",
     "subheading": "Supporting the Research Journey"},
    {"type": "bullets", "heading": "Our approach", "subheading": "One open stack",
     "bullets": ["Harmonise data to established standards", "Make collections discoverable",
                 "Federated analysis — data never moves"]},
    {"type": "stack", "heading": "Our approach", "subheading": "One open stack — same at every scale",
     "rows": [("Local data", "Harmonise to established standards"), ("Catalogue", "MOLGENIS meta-data catalogue"),
              ("Request", "BBMRI-ERIC Negotiator"), ("Access", "LS AAI + OIDC"),
              ("Analyse", "DataSHIELD + Molgenis Armadillo")]},
    {"type": "image", "side": "left", "heading": "Features", "subheading": "Manage projects and data",
     "bullets": ["Data are grouped into projects", "Each project can be shared with users",
                 "Admins manage projects from the UI"],
     "image": os.path.join(ASSET_ROOT, "molgenis-armadillo/public/ui-projects.png")},
    {"type": "photo", "side": "right", "heading": "The problem", "subheading": "A fragmented researcher journey",
     "bullets": ["Large quantities of existing data", "Combining sources boosts power",
                 "Fragmented landscape to reuse data"],
     "image": _p("public/screenshot-combining-data.jpeg")},
    {"type": "demo", "step": "Catalogue", "image": _p("public/screenshot-catalogue-search.png"),
     "text": "Search for birth cohorts with relevant harmonised data"},
    {"type": "demo", "step": "Analyse", "image": _p("public/screenshot-forest-plot.jpg"),
     "text": "Preliminary results: mode of delivery and postpartum depression"},
    {"type": "table", "heading": "Armadillo vs Opal", "subheading": "Two DataSHIELD servers",
     "rows": [["Feature", "Armadillo", "Opal"], ["Storage", "Parquet files", "Relational DB"],
              ["Deployment", "Lightweight", "Heavier"]]},
]

if __name__ == "__main__":
    out = "/tmp/molgenis-demo.pptx"
    DeckBuilder().build(DEMO_DECK).save(out)
    print("wrote", out)
